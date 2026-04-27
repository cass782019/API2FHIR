"""Gerador de apresentação PPTX — FHIR-Forge Arquitetura.

Uso:
    uv run python scripts/generate_pptx.py
Saída:
    docs/fhir-forge-arquitetura.pptx  (17 slides, 16:9, PT-BR)
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # type: ignore[attr-defined]
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Paleta ───────────────────────────────────────────────────────────────────
DARK_NAVY    = RGBColor(0x0D, 0x21, 0x37)
ACCENT_BLUE  = RGBColor(0x1E, 0x6F, 0xDB)
ACCENT_GREEN = RGBColor(0x13, 0xA1, 0x63)
LIGHT_BG     = RGBColor(0xF0, 0xF4, 0xFA)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT    = RGBColor(0x1A, 0x1A, 0x2E)
MID_GREY     = RGBColor(0x6B, 0x72, 0x80)
AMBER        = RGBColor(0xD9, 0x77, 0x06)
LIGHT_BLUE   = RGBColor(0xDB, 0xEA, 0xFB)
LIGHT_GREEN  = RGBColor(0xD1, 0xF5, 0xE4)

FOOTER_TEXT  = "FHIR-Forge v0.1.0  |  Abril 2026  |  Apache-2.0  |  Cassiano Moralles"


# ── Helpers ──────────────────────────────────────────────────────────────────

def _new_slide(prs: Presentation) -> object:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # fundo claro padrão
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    return slide


def _set_bg(slide: object, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _txb(slide: object, text: str, left: float, top: float,
          width: float, height: float, *, size: int = 14,
          bold: bool = False, color: RGBColor = DARK_TEXT,
          align: object = PP_ALIGN.LEFT, italic: bool = False,
          wrap: bool = True) -> object:
    txb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _box(slide: object, left: float, top: float, width: float, height: float,
         fill: RGBColor = LIGHT_BG, border: RGBColor = ACCENT_BLUE,
         text: str = "", size: int = 12, text_color: RGBColor = DARK_TEXT,
         bold: bool = False, rounded: bool = True) -> object:
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: F401
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE = 5, but add_shape uses PP_MEDIA_TYPE id
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1)
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = text_color
    return shape


def _title_bar(slide: object, title: str,
               bg: RGBColor = ACCENT_BLUE) -> None:
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(0.85)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"  {title}"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE


def _footer(slide: object) -> None:
    _txb(slide, FOOTER_TEXT, 0.2, 7.15, 12.9, 0.3,
         size=9, color=MID_GREY, align=PP_ALIGN.CENTER)


def _arrow_right(slide: object, x: float, y: float, length: float = 0.5) -> None:
    """Seta horizontal simples via text box '→'."""
    _txb(slide, "→", x, y, length, 0.35,
         size=18, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)


def _multiline(slide: object, lines: list[str], left: float, top: float,
               width: float, height: float, *, size: int = 13,
               color: RGBColor = DARK_TEXT, bold_first: bool = False) -> None:
    txb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (i == 0 and bold_first)


# ── Slides ───────────────────────────────────────────────────────────────────

def slide_01_capa(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _set_bg(slide, DARK_NAVY)

    # Barra verde no topo
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_GREEN
    bar.line.fill.background()

    _txb(slide, "FHIR-Forge", 0.7, 1.8, 11.9, 1.5,
         size=60, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _txb(slide,
         "Pipeline Híbrido LLM + SUSHI para Interoperabilidade em Saúde Digital\n"
         "Conversão Automática de APIs em Recursos FHIR R4 · BR Core · RNDS",
         0.7, 3.5, 11.5, 1.4,
         size=22, color=RGBColor(0xB0, 0xC8, 0xF0), align=PP_ALIGN.LEFT)

    # Barra inferior
    bot = slide.shapes.add_shape(1, Inches(0), Inches(6.9), Inches(13.33), Inches(0.6))
    bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT_GREEN
    bot.line.fill.background()
    _txb(slide,
         "Versão 0.1.0  |  Abril 2026  |  Apache-2.0  |  Cassiano Moralles",
         0, 6.9, 13.33, 0.6,
         size=13, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    # Badge canto superior direito
    _txb(slide, "FHIR R4  ·  BR Core  ·  RNDS",
         9.5, 0.25, 3.5, 0.4,
         size=11, color=RGBColor(0x80, 0xB0, 0xFF), align=PP_ALIGN.RIGHT)


def slide_02_agenda(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Agenda")

    itens = [
        "1.  Problema e Motivação",
        "2.  O que é FHIR, BR Core e RNDS",
        "3.  Arquitetura Geral do Sistema",
        "4.  Fluxo de Dados: OpenAPI → Bundle FHIR",
        "5.  Agente LangGraph: os 5 Nós",
        "6.  Roteamento de LLMs e Privacidade (LGPD)",
        "7.  Infraestrutura Docker",
        "8.  Segurança e Conformidade",
        "9.  Terminologia Clínica (SNOMED / LOINC)",
        "10. API REST e Workers Assíncronos",
        "11. Servidor MCP: Integração com Claude",
        "12. Pipeline de Avaliação de Qualidade",
        "13. Testes e Qualidade",
        "14. Mapa de Pacotes (uv Workspace)",
        "15. Resultados e Próximos Passos",
    ]
    col1 = itens[:8]
    col2 = itens[8:]

    _multiline(slide, col1, 0.5, 1.0, 6.2, 5.8, size=15, color=DARK_TEXT)
    _multiline(slide, col2, 6.9, 1.0, 6.2, 5.8, size=15, color=DARK_TEXT)
    _footer(slide)


def slide_03_problema(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "O Problema: APIs de Saúde vs. Interoperabilidade FHIR")

    pain = [
        ("Heterogeneidade",
         "Centenas de sistemas de saúde expõem APIs REST proprietárias — formatos "
         "inconsistentes, sem semântica clínica padronizada."),
        ("Exigência da RNDS",
         "A Rede Nacional de Dados em Saúde (RNDS) exige FHIR R4 + BR Core. "
         "A conversão manual é cara, lenta e sujeita a erros."),
        ("Necessidade de Automação",
         "Arquitetos precisam de uma camada de tradução automática e auditável "
         "entre o legado e o padrão nacional de interoperabilidade."),
    ]
    for i, (titulo, desc) in enumerate(pain):
        top = 1.1 + i * 1.7
        _box(slide, 0.4, top, 8.5, 1.45, fill=LIGHT_BLUE, border=ACCENT_BLUE)
        _txb(slide, titulo, 0.6, top + 0.05, 8.1, 0.4, size=14, bold=True, color=ACCENT_BLUE)
        _txb(slide, desc, 0.6, top + 0.45, 8.1, 0.9, size=12, color=DARK_TEXT)

    stats = [
        ("~4.000", "sistemas na RNDS"),
        ("LGPD", "dados sensíveis"),
        ("FHIR R4", "padrão do MS Brasil"),
    ]
    for i, (val, label) in enumerate(stats):
        top = 1.1 + i * 1.7
        _box(slide, 9.2, top, 3.9, 1.45, fill=ACCENT_GREEN, border=ACCENT_GREEN)
        _txb(slide, val, 9.2, top + 0.1, 3.9, 0.7,
             size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, label, 9.2, top + 0.8, 3.9, 0.5,
             size=13, color=WHITE, align=PP_ALIGN.CENTER)

    _footer(slide)


def slide_04_fhir(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "FHIR: O Padrão de Interoperabilidade em Saúde")

    cols = [
        ("FHIR R4", ACCENT_BLUE,
         ["Fast Healthcare Interoperability Resources",
          "• Formato JSON / XML / RDF",
          "• API RESTful padronizada",
          "• Recursos tipados: Patient,", "  Encounter, Observation…",
          "• Padrão HL7 Internacional"]),
        ("BR Core", ACCENT_GREEN,
         ["Perfil brasileiro do FHIR R4",
          "• Mantido pelo DATASUS / MS",
          "• CPF e CNS como identificadores",
          "• Terminologia SNOMED CT Edição BR",
          "• Versão 1.0.0 (2024)"]),
        ("RNDS", DARK_NAVY,
         ["Reg. Nacional de Dados em Saúde",
          "• Plataforma federal de troca",
          "  de dados clínicos",
          "• Autenticação mTLS + OAuth2",
          "• Ambientes de homologação",
          "  e produção"]),
    ]
    for i, (titulo, cor, linhas) in enumerate(cols):
        left = 0.4 + i * 4.3
        _box(slide, left, 1.0, 4.0, 0.6, fill=cor, border=cor)
        _txb(slide, titulo, left, 1.0, 4.0, 0.6,
             size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _multiline(slide, linhas, left + 0.1, 1.7, 3.8, 4.5, size=13, color=DARK_TEXT)

    _txb(slide,
         "✦  FHIR-Forge elimina a necessidade de escrever conversores FHIR manualmente.",
         0.5, 6.5, 12.3, 0.45,
         size=14, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    _footer(slide)


def slide_05_arquitetura(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Arquitetura do Sistema — Visão Geral")

    boxes = [
        (0.3,  "Entrada\nOpenAPI/Swagger\n(YAML ou JSON)", LIGHT_BLUE,   ACCENT_BLUE, DARK_TEXT),
        (2.9,  "SwaggerLens\nParser",                      ACCENT_BLUE,  ACCENT_BLUE, WHITE),
        (5.5,  "FhirForge\nAgente LangGraph",              ACCENT_GREEN, ACCENT_GREEN, WHITE),
        (8.1,  "HAPI FHIR\n$validate",                     DARK_NAVY,    DARK_NAVY,   WHITE),
        (10.7, "Saída\nBundle FHIR R4\nBR Core / RNDS",   ACCENT_GREEN, ACCENT_GREEN, WHITE),
    ]
    for left, txt, fill, border, tcol in boxes:
        _box(slide, left, 2.2, 2.3, 1.6, fill=fill, border=border, text=txt,
             size=12, text_color=tcol, bold=True)

    for x in [2.65, 5.25, 7.85, 10.45]:
        _arrow_right(slide, x, 2.8, 0.4)

    # Satélites (linha tracejada simulada com caixa menor)
    _box(slide, 3.8, 4.4, 2.5, 0.9, fill=LIGHT_BG, border=MID_GREY,
         text="TermMapper\nSNOMED · LOINC · CID-10", size=11, text_color=MID_GREY)
    _box(slide, 6.4, 4.4, 2.5, 0.9, fill=LIGHT_BG, border=MID_GREY,
         text="LLMs\nClaude Opus 4.7 / vLLM", size=11, text_color=MID_GREY)

    _txb(slide, "↑", 4.95, 3.85, 0.4, 0.5, size=16, color=MID_GREY, align=PP_ALIGN.CENTER)
    _txb(slide, "↑", 7.55, 3.85, 0.4, 0.5, size=16, color=MID_GREY, align=PP_ALIGN.CENTER)

    _txb(slide,
         "Pipeline stateful com checkpoints Postgres — cada etapa rastreada no Langfuse",
         0.5, 5.65, 12.3, 0.4,
         size=13, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)
    _footer(slide)


def slide_06_fluxo(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Fluxo de Dados: Da Especificação ao Recurso FHIR")

    steps = [
        ("1. Recepção",    ACCENT_BLUE,
         "Arquivo .yaml/.json recebido via POST /convert ou fila Dramatiq. Job ID gerado via ULID."),
        ("2. SwaggerLens", ACCENT_BLUE,
         "Valida o spec (openapi-spec-validator), extrai endpoints tipados: path, method, operationId, schemas, $refs resolvidos."),
        ("3. Mapeamento LLM", ACCENT_GREEN,
         "Claude Opus 4.7 recebe cada endpoint com exemplos few-shot BR Core → gera JSON do recurso FHIR estruturado."),
        ("4. Validação HAPI", DARK_NAVY,
         "POST /$validate no HAPI FHIR 8.4.0 contra perfil BR Core. Erros → fix_node (LLM corrige, até 3 tentativas)."),
        ("5. Bundle Final", ACCENT_GREEN,
         "Bundle FHIR R4 com todos os recursos válidos. Entregue ao cliente ou submetido à RNDS via mTLS."),
    ]
    for i, (titulo, cor, desc) in enumerate(steps):
        top = 1.0 + i * 1.15
        _box(slide, 0.3, top, 2.2, 0.9, fill=cor, border=cor)
        _txb(slide, titulo, 0.3, top + 0.1, 2.2, 0.7,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, desc, 2.7, top + 0.05, 10.0, 0.85, size=12, color=DARK_TEXT)
        if i < 4:
            _txb(slide, "↓", 1.1, top + 0.9, 0.8, 0.25,
                 size=14, color=cor, align=PP_ALIGN.CENTER)

    _footer(slide)


def slide_07_langgraph(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Agente de Conversão: Grafo LangGraph com 5 Nós")

    nodes = [
        (1.0,  2.0, "parse_node",      ACCENT_BLUE,  WHITE),
        (4.0,  2.0, "mapping_node",    ACCENT_BLUE,  WHITE),
        (7.0,  2.0, "validation_node", ACCENT_BLUE,  WHITE),
        (10.0, 2.0, "bundle_node",     ACCENT_GREEN, WHITE),
        (7.0,  4.2, "fix_node",        AMBER,        WHITE),
    ]
    for left, top, nome, fill, tc in nodes:
        _box(slide, left, top, 2.5, 0.75, fill=fill, border=fill, text=nome, size=12, text_color=tc, bold=True)

    # Setas principais
    for x in [3.55, 6.55, 9.55]:
        _arrow_right(slide, x, 2.25, 0.4)

    # Seta validation → fix
    _txb(slide, "↓ inválido\n(retry < 3)", 7.0, 2.8, 2.5, 0.7, size=10, color=AMBER, align=PP_ALIGN.CENTER)
    # Seta fix → validation (retorno)
    _txb(slide, "↑ retry", 9.65, 3.8, 1.0, 0.35, size=10, color=AMBER, align=PP_ALIGN.CENTER)
    _txb(slide, "→", 10.55, 4.42, 0.5, 0.35, size=14, color=AMBER, align=PP_ALIGN.CENTER)
    # Seta validation → bundle (válido ou retry≥3)
    _txb(slide, "✓ válido /\nretry ≥ 3 →", 9.6, 2.05, 1.2, 0.7, size=10, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    # Descrições dos nós
    descs = [
        (0.3,  5.2, "parse_node",       "Extrai operation_ids dos endpoints a processar"),
        (3.3,  5.2, "mapping_node",     "LLM few-shot BR Core → dicts FHIR"),
        (6.3,  5.2, "validation_node",  "POST $validate no HAPI; popula validation_errors"),
        (9.3,  5.2, "bundle_node",      "Monta Bundle R4 com meta.tag fhir-forge"),
        (6.3,  5.9, "fix_node",         "LLM corrige recurso com erros como contexto"),
    ]
    for left, top, nome, desc in descs:
        _txb(slide, f"• {nome}: {desc}", left, top, 3.7, 0.55, size=10, color=DARK_TEXT)

    _txb(slide,
         "Estado persiste no Postgres via AsyncPostgresSaver — retomada automática em falha",
         0.5, 6.6, 12.3, 0.35,
         size=11, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)
    _footer(slide)


def slide_08_llm_lgpd(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Roteamento de LLMs: Privacidade por Design (LGPD)")

    # Árvore de decisão
    _box(slide, 0.4, 1.0, 3.5, 0.7, fill=LIGHT_BLUE, border=ACCENT_BLUE,
         text="Dados de entrada", size=13, text_color=DARK_TEXT, bold=True)
    _txb(slide, "↓", 1.9, 1.75, 0.8, 0.4, size=18, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    _box(slide, 0.4, 2.2, 3.5, 0.7, fill=AMBER, border=AMBER,
         text="Contém PHI?", size=14, text_color=WHITE, bold=True)

    _txb(slide, "NÃO →", 4.1, 2.35, 1.2, 0.4, size=13, color=ACCENT_GREEN, bold=True)
    _txb(slide, "SIM + egress=false →", 0.0, 3.05, 4.5, 0.4, size=13, color=AMBER, bold=True)

    nao_routes = [
        ("claude-opus-4-7", "Anthropic (padrão)"),
        ("claude-haiku-4-5-20251001", "Anthropic (rápido)"),
        ("claude-sonnet-4-6", "AWS Bedrock (fallback)"),
        ("sabia-4", "Maritaca (PT-BR nativo)"),
    ]
    for i, (model, provider) in enumerate(nao_routes):
        top = 1.85 + i * 0.62
        _box(slide, 5.4, top, 7.5, 0.52, fill=LIGHT_GREEN, border=ACCENT_GREEN,
             text=f"{model}  —  {provider}", size=12, text_color=DARK_TEXT)

    _box(slide, 0.4, 3.5, 12.5, 0.65, fill=AMBER, border=AMBER,
         text="vLLM on-prem (Llama 3.3 70B)  —  dado NUNCA sai da infraestrutura",
         size=13, text_color=WHITE, bold=True)

    _txb(slide,
         "⚠  FEATURE_ALLOW_PHI_EGRESS=false por padrão. "
         "PHI = dado pessoal de saúde protegido pela LGPD (Lei 13.709/2018).",
         0.4, 4.45, 12.5, 0.55,
         size=12, bold=True, color=AMBER)

    _txb(slide,
         "SecretStr (pydantic): chaves de API nunca aparecem em logs ou traces Langfuse.",
         0.4, 5.1, 12.5, 0.4, size=12, color=DARK_TEXT, italic=True)

    _footer(slide)


def slide_09_infra(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Infraestrutura: Stack Docker Completa")

    services = [
        ("HAPI FHIR 8.4.0", "Porta 8090", "Validação + armazenamento FHIR", ACCENT_BLUE),
        ("Snowstorm 7.5.0", "Porta 8080", "Servidor SNOMED CT", ACCENT_BLUE),
        ("Elasticsearch 8.11.1", "Porta 9200", "Backend do Snowstorm", ACCENT_BLUE),
        ("PostgreSQL 16 + pgvector", "Porta 5432", "HAPI JPA · LangGraph · App data", DARK_NAVY),
        ("Redis 7.4", "Porta 6379", "Fila Dramatiq + cache de terminologia", DARK_NAVY),
        ("MinIO", "Portas 9000/9001", "Artefatos S3-compatible", DARK_NAVY),
        ("Langfuse 3.x", "Porta 3000", "Observabilidade de LLMs (traces)", ACCENT_GREEN),
        ("FastAPI App", "Porta 8000", "Gateway REST da aplicação", ACCENT_GREEN),
        ("Dramatiq Workers", "—", "Processamento assíncrono de conversões", ACCENT_GREEN),
    ]
    for i, (nome, porta, desc, cor) in enumerate(services):
        col = i % 3
        row = i // 3
        left = 0.3 + col * 4.35
        top = 1.0 + row * 1.95
        _box(slide, left, top, 4.1, 0.5, fill=cor, border=cor)
        _txb(slide, nome, left, top + 0.05, 4.1, 0.45,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, porta, left, top + 0.55, 4.1, 0.35,
             size=11, color=MID_GREY, align=PP_ALIGN.CENTER)
        _txb(slide, desc, left, top + 0.9, 4.1, 0.65,
             size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    _txb(slide,
         "Todos os serviços com healthchecks configurados.  `make up` sobe toda a stack.  ~12 GB RAM necessários.",
         0.3, 6.9, 12.7, 0.35,
         size=11, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)
    _footer(slide)


def slide_10_seguranca(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Segurança e Conformidade: LGPD + RNDS")

    pillars = [
        ("Autenticação", ACCENT_BLUE,
         ["• Bearer token na API REST",
          "• Feature flag ADMIN_NOAUTH (dev)",
          "• OAuth2 + mTLS para RNDS",
          "• Tokens com cache e TTL"]),
        ("PHI / LGPD", AMBER,
         ["• FEATURE_ALLOW_PHI_EGRESS=false",
          "• LLM routing auto p/ on-prem",
          "• Sem PHI em logs ou traces",
          "• Auditabilidade via Langfuse"]),
        ("Credenciais", DARK_NAVY,
         ["• SecretStr pydantic p/ secrets",
          "• Certs RNDS em secrets/ (não commitado)",
          "• Variáveis de ambiente no .env",
          "• Langfuse keys separadas"]),
        ("Resiliência", ACCENT_GREEN,
         ["• tenacity retry (3x backoff exp.)",
          "• pybreaker circuit breaker",
          "• Dramatiq DLQ após 3 falhas",
          "• Nenhuma tarefa perdida silenciosa"]),
    ]
    for i, (titulo, cor, items) in enumerate(pillars):
        left = 0.3 + i * 3.2
        _box(slide, left, 1.0, 3.0, 0.55, fill=cor, border=cor)
        _txb(slide, titulo, left, 1.05, 3.0, 0.45,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _multiline(slide, items, left + 0.05, 1.65, 2.9, 3.5, size=12, color=DARK_TEXT)

    _txb(slide,
         "Conformidade LGPD verificada em design — dado sensível jamais egresso sem autorização explícita (Art. 11, Lei 13.709/2018).",
         0.3, 6.5, 12.7, 0.45,
         size=12, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    _footer(slide)


def slide_11_terminologia(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Terminologia Clínica: SNOMED CT, LOINC, CID-10, TUSS")

    # Diagrama de camadas
    camadas = [
        (0.4, 1.0, 5.5, 0.85, ACCENT_BLUE,  WHITE,  "HAPI FHIR  —  $expand, $lookup (fallback)"),
        (0.4, 1.95, 5.5, 0.85, ACCENT_GREEN, WHITE, "Snowstorm  —  find_concept (ECL), $translate"),
        (0.4, 2.9, 5.5, 0.85, DARK_NAVY,    WHITE,  "Redis Cache  —  TTL 3600s, chave {system}:{code}"),
    ]
    _txb(slide, "Camadas de Lookup (ordem de prioridade: ↑ da base para o topo)",
         0.4, 0.9, 5.7, 0.4, size=11, italic=True, color=MID_GREY)
    for left, top, w, h, fill, tc, txt in camadas:
        _box(slide, left, top, w, h, fill=fill, border=fill, text=txt, size=13, text_color=tc, bold=True)

    # Tabela de sistemas
    headers_s = ["Sistema", "Uso Principal"]
    rows_s = [
        ["SNOMED CT BR", "Diagnósticos, procedimentos clínicos"],
        ["LOINC", "Observações e exames laboratoriais"],
        ["CID-10", "Condições (sistema nacional, DATASUS)"],
        ["TUSS", "Procedimentos da saúde suplementar (ANS)"],
        ["SIGTAP", "Procedimentos do SUS (DATASUS)"],
    ]
    _box(slide, 6.2, 1.0, 7.0, 0.5, fill=DARK_NAVY, border=DARK_NAVY,
         text=f"{headers_s[0]}              {headers_s[1]}", size=13, text_color=WHITE, bold=True)
    for i, row in enumerate(rows_s):
        fill_r = LIGHT_BG if i % 2 == 0 else WHITE
        _box(slide, 6.2, 1.55 + i * 0.55, 7.0, 0.5, fill=fill_r, border=LIGHT_BG)
        _txb(slide, row[0], 6.3, 1.57 + i * 0.55, 2.5, 0.45, size=12, bold=True, color=DARK_TEXT)
        _txb(slide, row[1], 8.9, 1.57 + i * 0.55, 4.2, 0.45, size=12, color=DARK_TEXT)

    _txb(slide,
         "Lookup em cache Redis primeiro — 0 chamadas HTTP em cache hits. "
         "Mapper orquestra: Redis → Snowstorm → HAPI fallback.",
         0.4, 5.5, 12.7, 0.55, size=12, italic=True, color=MID_GREY)
    _footer(slide)


def slide_12_api_workers(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Interface REST + Processamento Assíncrono")

    # FastAPI
    _txb(slide, "API REST — FastAPI (Porta 8000)",
         0.3, 1.0, 6.2, 0.45, size=15, bold=True, color=ACCENT_BLUE)
    endpoints = [
        ("GET  /health", "Status de todos os serviços (HAPI, Redis, Postgres)"),
        ("POST /convert", "Conversão síncrona ou assíncrona (Dramatiq)"),
        ("GET  /fhir/{type}/{id}", "Proxy para HAPI FHIR"),
        ("POST /fhir/{type}/$validate", "Proxy de validação HAPI"),
        ("*    /mcp", "Ferramentas MCP para Claude (se habilitado)"),
    ]
    for i, (ep, desc) in enumerate(endpoints):
        top = 1.55 + i * 0.85
        _box(slide, 0.3, top, 6.2, 0.75, fill=LIGHT_BLUE, border=ACCENT_BLUE)
        _txb(slide, ep, 0.45, top + 0.05, 2.8, 0.35, size=11, bold=True, color=ACCENT_BLUE)
        _txb(slide, desc, 0.45, top + 0.38, 5.9, 0.35, size=11, color=DARK_TEXT)

    # Dramatiq
    _txb(slide, "Workers Assíncronos — Dramatiq + Redis",
         6.8, 1.0, 6.2, 0.45, size=15, bold=True, color=DARK_NAVY)
    worker_steps = [
        ("POST /convert (async_mode=true)", LIGHT_BLUE, ACCENT_BLUE),
        ("Redis Queue \"fhir\"",             LIGHT_BG,   DARK_NAVY),
        ("convert_spec_actor",              ACCENT_GREEN, ACCENT_GREEN),
        ("Resultado em Redis (job_id)",     LIGHT_BG,   DARK_NAVY),
        ("DLQ após 3 falhas → dlq_handler", LIGHT_BG,  AMBER),
    ]
    for i, (txt, fill_c, border_c) in enumerate(worker_steps):
        top = 1.55 + i * 0.9
        _box(slide, 6.8, top, 6.2, 0.75, fill=fill_c, border=border_c, text=txt, size=12, text_color=DARK_TEXT)
        if i < 4:
            _txb(slide, "↓", 9.6, top + 0.75, 0.5, 0.2,
                 size=12, color=DARK_NAVY, align=PP_ALIGN.CENTER)

    _txb(slide,
         "OTEL spans automáticos (FastAPIInstrumentor).  structlog com request_id em cada request.",
         0.3, 6.9, 12.7, 0.35, size=11, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)
    _footer(slide)


def slide_13_mcp(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Servidor MCP: 6 Ferramentas FHIR para Claude")

    _txb(slide, "Claude (IA)", 0.3, 1.1, 2.5, 0.65,
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    shape = slide.shapes.add_shape(1, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.65))
    shape.fill.solid(); shape.fill.fore_color.rgb = DARK_NAVY
    shape.line.fill.background()

    _txb(slide, "⇄", 2.95, 1.22, 0.6, 0.45, size=22, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    _box(slide, 3.6, 1.0, 3.2, 0.85, fill=ACCENT_BLUE, border=ACCENT_BLUE,
         text="Servidor MCP\n(FastMCP 3.x / /mcp)", size=12, text_color=WHITE, bold=True)

    _txb(slide, "⇄", 6.9, 1.22, 0.6, 0.45, size=22, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    _box(slide, 7.55, 0.95, 2.5, 0.95, fill=LIGHT_BLUE, border=ACCENT_BLUE,
         text="HAPI FHIR\n8090", size=12, text_color=DARK_TEXT, bold=True)
    _box(slide, 10.2, 0.95, 2.9, 0.95, fill=LIGHT_GREEN, border=ACCENT_GREEN,
         text="Snowstorm\n8080", size=12, text_color=DARK_TEXT, bold=True)

    tools = [
        ("validate_resource", "Valida recurso FHIR contra perfil BR Core no HAPI"),
        ("get_resource", "Busca recurso FHIR por tipo + ID no HAPI"),
        ("search_resources", "Busca FHIR com parâmetros query-string, paginação automática"),
        ("expand_valueset", "Expande ValueSet FHIR via $expand"),
        ("find_concept", "Busca conceito SNOMED CT por texto ou ECL"),
        ("translate_code", "Traduz código entre sistemas via ConceptMap/$translate"),
    ]
    _box(slide, 0.3, 2.15, 4.0, 0.45, fill=DARK_NAVY, border=DARK_NAVY,
         text="Ferramenta", size=13, text_color=WHITE, bold=True)
    _box(slide, 4.4, 2.15, 8.6, 0.45, fill=DARK_NAVY, border=DARK_NAVY,
         text="Descrição", size=13, text_color=WHITE, bold=True)
    for i, (nome, desc) in enumerate(tools):
        fill_r = LIGHT_BG if i % 2 == 0 else WHITE
        _box(slide, 0.3, 2.65 + i * 0.65, 4.0, 0.6, fill=fill_r, border=LIGHT_BG)
        _txb(slide, nome, 0.4, 2.67 + i * 0.65, 3.8, 0.55,
             size=12, bold=True, color=ACCENT_BLUE)
        _box(slide, 4.4, 2.65 + i * 0.65, 8.6, 0.6, fill=fill_r, border=LIGHT_BG)
        _txb(slide, desc, 4.5, 2.67 + i * 0.65, 8.4, 0.55, size=12, color=DARK_TEXT)

    _txb(slide,
         "Montado em /mcp via FastMCP 3.x (streamable-http). Ativo quando FEATURE_MCP_ENABLED=true.",
         0.3, 6.85, 12.7, 0.35, size=11, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)
    _footer(slide)


def slide_14_eval(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Avaliação de Qualidade: Precision, Recall e Golden Files")

    # Fórmulas
    formulas = [
        ("Precision", "| Tipos preditos ∩ esperados |  /  | Tipos preditos |",
         "Mede: quantos recursos gerados estão corretos?"),
        ("Recall", "| Tipos preditos ∩ esperados |  /  | Tipos esperados |",
         "Mede: quantos recursos esperados foram gerados?"),
        ("Semantic Diff", "deepdiff(predito, esperado, ignore_order=True)",
         "Ignora: id, meta.lastUpdated, meta.versionId (campos voláteis)"),
    ]
    for i, (titulo, formula, nota) in enumerate(formulas):
        top = 1.0 + i * 1.5
        _box(slide, 0.3, top, 3.0, 1.3, fill=ACCENT_BLUE, border=ACCENT_BLUE)
        _txb(slide, titulo, 0.3, top + 0.1, 3.0, 0.55,
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, formula, 3.5, top + 0.05, 9.5, 0.55, size=13, bold=True, color=DARK_TEXT)
        _txb(slide, nota, 3.5, top + 0.65, 9.5, 0.5, size=12, italic=True, color=MID_GREY)

    # EvalRunner flow
    _txb(slide, "EvalRunner:", 0.3, 5.55, 3.5, 0.4, size=13, bold=True, color=DARK_NAVY)
    flow = ["data/golden/swagger_specs/*.yaml", "→  conversão (mock ou real)", "→  EvalReport(pairs=[PairResult(...)])"]
    for i, txt in enumerate(flow):
        _txb(slide, txt, 3.9 + i * 3.0, 5.55, 2.8, 0.4, size=12, color=DARK_TEXT)
        if i < 2:
            _txb(slide, "→", 6.7 + i * 3.0, 5.6, 0.4, 0.3, size=13, color=ACCENT_GREEN)

    _box(slide, 0.3, 6.2, 12.7, 0.55, fill=LIGHT_GREEN, border=ACCENT_GREEN,
         text="Regressão bloqueante: diff semântico não-zero em golden = bloqueador de merge  |  Precision baseline ≥ 0.80",
         size=12, text_color=DARK_TEXT, bold=True)
    _footer(slide)


def slide_15_testes(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Estratégia de Testes: 252 Testes, 91% de Cobertura")

    headers = ["Camada", "Marker", "Ferramentas", "O que testa", "Cobertura mín."]
    rows = [
        ["Unit", "@pytest.mark.unit", "pytest, respx, unittest.mock", "Lógica pura — 0 I/O real", "≥ 80%"],
        ["Integration", "@pytest.mark.integration", "testcontainers (Postgres+Redis+HAPI)", "CRUD HAPI, filas, banco", "≥ 80%"],
        ["Regression", "@pytest.mark.regression", "deepdiff golden-files", "Output FHIR não regride", "—"],
        ["E2E", "@pytest.mark.e2e", "Docker stack completa", "Saga OpenAPI → Bundle", "—"],
        ["RNDS", "@pytest.mark.rnds", "Sandbox RNDS", "Submissão real (credenciais)", "—"],
    ]
    cols_w = [1.6, 2.4, 3.8, 3.5, 1.6]
    lefts = [0.25]
    for w in cols_w[:-1]:
        lefts.append(lefts[-1] + w + 0.05)

    for j, (h, lft) in enumerate(zip(headers, lefts)):
        _box(slide, lft, 1.0, cols_w[j], 0.5, fill=DARK_NAVY, border=DARK_NAVY,
             text=h, size=12, text_color=WHITE, bold=True)

    fills = [LIGHT_BG, WHITE, LIGHT_BG, WHITE, LIGHT_BG]
    for i, row in enumerate(rows):
        top = 1.55 + i * 0.75
        for j, (cell, lft) in enumerate(zip(row, lefts)):
            _box(slide, lft, top, cols_w[j], 0.7, fill=fills[i], border=LIGHT_BG)
            _txb(slide, cell, lft + 0.05, top + 0.08, cols_w[j] - 0.1, 0.58,
                 size=10, color=DARK_TEXT)

    # KPIs
    kpis = [("252", "testes unit"), ("91%", "cobertura"), ("0", "erros ruff"), ("0", "erros mypy")]
    for i, (val, label) in enumerate(kpis):
        left = 0.25 + i * 3.3
        _box(slide, left, 5.55, 3.1, 0.9, fill=ACCENT_GREEN, border=ACCENT_GREEN)
        _txb(slide, val, left, 5.6, 3.1, 0.5, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, label, left, 6.1, 3.1, 0.3, size=11, color=WHITE, align=PP_ALIGN.CENTER)

    _footer(slide)


def slide_16_packages(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Organização: uv Workspace — 7 Packages + 2 Apps")

    # Grafo de dependências (texto visual)
    deps = [
        (0.4, 1.0, "packages/core", DARK_NAVY, WHITE,
         "Settings · LLMRouter · types · exceptions · logging\n(sem dependências internas — é a base de tudo)"),
        (0.4, 2.3, "packages/swagger_lens", ACCENT_BLUE, WHITE, "→ core  |  Parser OpenAPI 2.0/3.x"),
        (0.4, 3.1, "packages/fhir_forge", ACCENT_BLUE, WHITE, "→ core · swagger_lens · connectors  |  LangGraph 5 nós"),
        (0.4, 3.9, "packages/term_mapper", ACCENT_BLUE, WHITE, "→ core  |  SNOMED/LOINC/CID-10/TUSS + Redis cache"),
        (0.4, 4.7, "packages/connectors", ACCENT_BLUE, WHITE, "→ core  |  HAPI client + RNDS client (mTLS)"),
        (0.4, 5.5, "packages/eval", ACCENT_BLUE, WHITE, "→ core · fhir_forge  |  EvalRunner + métricas"),
        (0.4, 6.3, "packages/mcp_server", ACCENT_BLUE, WHITE, "→ core  |  6 FastMCP tools FHIR"),
    ]
    for left, top, nome, fill_c, tc, desc in deps:
        _box(slide, left, top, 3.5, 0.6, fill=fill_c, border=fill_c, text=nome, size=11, text_color=tc, bold=True)
        _txb(slide, desc, 4.1, top + 0.08, 4.5, 0.5, size=11, color=DARK_TEXT)

    # Apps
    apps = [
        (8.9, 2.3, "apps/api", ACCENT_GREEN,
         "→ fhir_forge · connectors · mcp_server · core\nFastAPI gateway — 4 routers + OTEL"),
        (8.9, 3.6, "apps/worker", ACCENT_GREEN,
         "→ fhir_forge · core\nDramatiq actors + DLQ handler"),
    ]
    _txb(slide, "Apps:", 8.9, 1.75, 4.0, 0.45, size=14, bold=True, color=DARK_NAVY)
    for left, top, nome, fill_c, desc in apps:
        _box(slide, left, top, 3.5, 0.6, fill=fill_c, border=fill_c, text=nome, size=12, text_color=WHITE, bold=True)
        _txb(slide, desc, 12.55, top + 0.0, 0.5, 0.65, size=10, color=DARK_TEXT)
        _multiline(slide, desc.split("\n"), 8.9, top + 0.65, 4.3, 0.7, size=11, color=DARK_TEXT)

    _txb(slide,
         "Cada membro: pyproject.toml próprio + src layout + hatchling build.  ruff + mypy rodados na raiz.",
         0.3, 7.15, 12.7, 0.3, size=10, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)
    _footer(slide)


def slide_17_resultados(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _title_bar(slide, "Resultados Atuais e Roadmap")

    kpis = [
        ("252", "testes unit passando"),
        ("91%", "cobertura de código"),
        ("9/9", "fases completas"),
        ("0", "erros lint/mypy"),
    ]
    for i, (val, label) in enumerate(kpis):
        left = 0.3 + i * 3.2
        _box(slide, left, 1.0, 3.0, 1.6, fill=ACCENT_BLUE, border=ACCENT_BLUE)
        _txb(slide, val, left, 1.05, 3.0, 1.0, size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, label, left, 2.0, 3.0, 0.55, size=13, color=WHITE, align=PP_ALIGN.CENTER)

    _txb(slide, "Próximos Passos:", 0.3, 3.0, 12.7, 0.45, size=15, bold=True, color=DARK_NAVY)
    roadmap = [
        "• GitHub Actions CI/CD (ci.yml, ig-build.yml) — push/PR/main em < 3/12/20 min",
        "• Helm charts e manifests Kubernetes para deploy em produção",
        "• IG FHIR (SUSHI + IG Publisher) com perfis BR Core customizados",
        "• Carga SNOMED CT Edição BR no Snowstorm (RF2 release)",
        "• Migrations Alembic — modelo de dados de persistência de jobs",
        "• Frontend de monitoramento de jobs (React + WebSocket)",
    ]
    _multiline(slide, roadmap, 0.5, 3.55, 12.5, 3.5, size=13, color=DARK_TEXT)

    # Barra final
    bot = slide.shapes.add_shape(1, Inches(0), Inches(6.9), Inches(13.33), Inches(0.6))
    bot.fill.solid(); bot.fill.fore_color.rgb = DARK_NAVY
    bot.line.fill.background()
    _txb(slide,
         "Apache-2.0  |  Autor: Cassiano Moralles  |  Python ≥ 3.12  |  FHIR R4 + BR Core 1.0.0",
         0, 6.92, 13.33, 0.55, size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_01_capa,
        slide_02_agenda,
        slide_03_problema,
        slide_04_fhir,
        slide_05_arquitetura,
        slide_06_fluxo,
        slide_07_langgraph,
        slide_08_llm_lgpd,
        slide_09_infra,
        slide_10_seguranca,
        slide_11_terminologia,
        slide_12_api_workers,
        slide_13_mcp,
        slide_14_eval,
        slide_15_testes,
        slide_16_packages,
        slide_17_resultados,
    ]
    for fn in builders:
        fn(prs)

    Path("docs").mkdir(exist_ok=True)
    out = Path("docs") / "fhir-forge-arquitetura.pptx"
    prs.save(str(out))
    print(f"Gerado: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
